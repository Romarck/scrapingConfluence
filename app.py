import streamlit as st
import os
import json
import requests
from slugify import slugify
from uuid import uuid4
from bs4 import BeautifulSoup
import base64
import datetime
import zipfile
import io  # For in-memory zip file creation
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation
from PIL import Image
from fpdf import FPDF  # For PDF output
import markdown  # For Markdown output

# Streamlit UI
st.title("Confluence Loader")

# Sidebar for configurations
st.sidebar.header("Confluence Configuration")
base_url = st.sidebar.text_input("Base URL", os.environ.get("CONFLUENCE_BASE_URL", ""))
space_key = st.sidebar.text_input("Space Key", os.environ.get("CONFLUENCE_SPACE_KEY", ""))
username = st.sidebar.text_input("Username", os.environ.get("CONFLUENCE_USERNAME", ""))
access_token = st.sidebar.text_input("Access Token", os.environ.get("CONFLUENCE_ACCESS_TOKEN", ""), type="password")
personal_access_token = st.sidebar.text_input("Personal Access Token", os.environ.get("CONFLUENCE_PERSONAL_ACCESS_TOKEN", ""), type="password")
cloud = st.sidebar.checkbox("Cloud", value=True)

# Output Format Configuration
st.sidebar.header("Output Format")
output_format = st.sidebar.selectbox("Select Output Format", ["txt", "md", "pdf", "json"])

# Attachment Handling Configuration
st.sidebar.header("Attachment Handling")
process_attachments = st.sidebar.checkbox("Process Attachments", value=True, help="Enable to extract text from attached files.")


# Functions
def valid_base_url(base_url):
    try:
        result = requests.get(base_url)
        result.raise_for_status()
        return True
    except requests.exceptions.RequestException as e:
        st.error(f"Invalid Base URL: {e}")
        return False

def fetch_confluence_data(url, username, access_token, personal_access_token):
    headers = {
        "Content-Type": "application/json",
        "Accept": "application/json",
    }
    if personal_access_token:
        headers["Authorization"] = f"Bearer {personal_access_token}"
    elif username and access_token:
        auth_string = f"{username}:{access_token}"
        auth_bytes = auth_string.encode('utf-8')
        auth_base64_bytes = base64.b64encode(auth_bytes)
        auth_base64_string = auth_base64_bytes.decode('utf-8')
        headers["Authorization"] = f"Basic {auth_base64_string}"
    else:
        st.error("Please provide either a personal access token or username and access token.")
        return None

    try:
        response = requests.get(url, headers=headers)
        response.raise_for_status()
        return response.json()
    except requests.exceptions.RequestException as e:
        st.error(f"Failed to fetch {url} from Confluence: {e}")
        return None

def fetch_all_pages_in_space(base_url, space_key, username, access_token, personal_access_token, cloud=True, start=0, limit=25, expand="body.storage,version,children.attachment"):  #Added children.attachment to expand
    try:
        confluence_url = base_url + ("/wiki" if cloud else "") + "/rest/api/content"
        url = f"{confluence_url}?spaceKey={space_key}&limit={limit}&start={start}&expand={expand}"
        data = fetch_confluence_data(url, username, access_token, personal_access_token)

        if data is None:
            return []

        if data.get("size") == 0:
            return []

        next_page_start = start + data["size"]
        next_page_results = fetch_all_pages_in_space(base_url, space_key, username, access_token, personal_access_token, cloud, next_page_start, limit, expand)
        return data["results"] + next_page_results
    except Exception as e:
        st.error(f"An error occurred while fetching pages: {e}")
        return []

def extract_code_blocks(content):
    soup = BeautifulSoup(content, 'html.parser')
    code_blocks = []
    for code_macro in soup.find_all('ac:structured-macro', attrs={'ac:name': 'code'}):
        language = code_macro.find('ac:parameter', attrs={'ac:name': 'language'})
        language = language.text if language else ""
        code_body = code_macro.find('ac:plain-text-body')
        if code_body:
            code = code_body.text
            code_blocks.append(f"\n```{language}\n{code.strip()}\n```\n")
    return ''.join(code_blocks)


def extract_text_from_attachment(attachment, username, access_token, personal_access_token, output_dir, base_url, cloud):  # Added base_url and cloud
    """Downloads an attachment and attempts to extract text content based on file type."""
    download_url = attachment['_links']['download']

    # Construct the correct URL, including /wiki for cloud instances
    confluence_url = base_url + ("/wiki" if cloud else "")
    full_download_url = confluence_url + download_url if not download_url.startswith('http') else download_url

    headers = {
        "Accept": "application/json",
    }
    if personal_access_token:
        headers["Authorization"] = f"Bearer {personal_access_token}"
    elif username and access_token:
        auth_string = f"{username}:{access_token}"
        auth_bytes = auth_string.encode('utf-8')
        auth_base64_bytes = base64.b64encode(auth_bytes)
        auth_base64_string = auth_base64_bytes.decode('utf-8')
        headers["Authorization"] = f"Basic {auth_base64_string}"
    else:
        st.error("Please provide either a personal access token or username and access token.")
        return None, None  # Return None for both text and download URL

    try:
        # Debugging: Print the URL before making the request
        print(f"Tentando baixar: {full_download_url}")  # Print the URL

        response = requests.get(full_download_url, headers=headers, stream=True)
        response.raise_for_status()

        file_extension = attachment['title'].split('.')[-1].lower()

        # Save the attachment to the "anexos" directory
        attachment_filepath = os.path.join(output_dir, "anexos", attachment['title'])  # Corrected path
        with open(attachment_filepath, "wb") as f:
            for chunk in response.iter_content(chunk_size=8192):  # Stream in chunks
                f.write(chunk)
        print(f"Anexo salvo com sucesso: {attachment['title']}")

        text = None  # Initialize text to None
        if file_extension == 'pdf':
            try:
                with open(attachment_filepath, 'rb') as f:  # Open saved file
                    pdf_reader = PyPDF2.PdfReader(f)
                    text = ""
                    for page_num in range(len(pdf_reader.pages)):
                        page = pdf_reader.pages[page_num]
                        text += page.extract_text()
            except Exception as e:
                text = f"Erro ao extrair texto do PDF {attachment['title']}: {e}"

        elif file_extension in ['doc', 'docx']:
            try:
                with open(attachment_filepath, 'rb') as f:  # Open saved file
                    document = Document(f)
                    text = '\n'.join([paragraph.text for paragraph in document.paragraphs])
            except Exception as e:
                text = f"Erro ao extrair texto do DOC/DOCX {attachment['title']}: {e}"

        elif file_extension in ['xls', 'xlsx']:
            try:
                with open(attachment_filepath, 'rb') as f:
                    workbook = openpyxl.load_workbook(f)
                    text = ""
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        for row in sheet.iter_rows():
                            row_values = [str(cell.value) for cell in row if cell.value is not None]
                            text += ', '.join(row_values) + '\n'
            except Exception as e:
                text = f"Erro ao extrair texto do XLS/XLSX {attachment['title']}: {e}"

        elif file_extension in ['ppt', 'pptx']:
            try:
                with open(attachment_filepath, 'rb') as f:
                    presentation = Presentation(f)
                    text = ""
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + '\n'
            except Exception as e:
                text = f"Erro ao extrair texto do PPT/PPTX {attachment['title']}: {e}"

        elif file_extension in ['png', 'jpg', 'jpeg']:
            try:
                image = Image.open(attachment_filepath)
                text = f"Imagem: {attachment['title']} (Não é possível extrair texto diretamente)"  # Just indicate it's an image

            except Exception as e:
                text = f"Erro ao processar imagem {attachment['title']}: {e}"
        else:
            text = f"Tipo de arquivo não suportado: {attachment['title']}"

        return text, full_download_url  # Return both text and download URL

    except requests.exceptions.RequestException as e:
        st.error(f"Falha ao baixar o anexo {attachment['title']}: {e}")
        return None, None  # Return None and None


def create_document_from_page(page, base_url, space_key, username, access_token, personal_access_token, cloud, process_attachments, output_dir):  # Added output_dir and base_url
    content = page.get("body", {}).get("storage", {}).get("value", "")
    if not content:
        content = ""

    code_blocks_text = extract_code_blocks(content)
    plain_text_content = BeautifulSoup(content, "html.parser").get_text()

    combined_text = plain_text_content + code_blocks_text
    text_with_preserved_structure = combined_text.replace('\n' * 3, '\n\n')

    if cloud:
        page_url = f"{base_url}/wiki/spaces/{space_key}/pages/{page.get('id')}"
    else:
        page_url = f"{base_url}/spaces/{space_key}/pages/{page.get('id')}"

    metadata = {
        "id": page.get("id"),
        "status": page.get("status"),
        "title": page.get("title"),
        "type": page.get("type"),
        "url": page_url,
        "version": page.get("version", {}).get("number"),
        "updated_by": page.get("version", {}).get("by", {}).get("displayName"),
        "updated_at": page.get("version", {}).get("when"),
    }

    attachment_text = ""
    if process_attachments:
        attachments = page.get('children', {}).get('attachment', {}).get('results', [])
        for attachment in attachments:
            attachment_content, download_url = extract_text_from_attachment(attachment, username, access_token, personal_access_token, output_dir, base_url, cloud)  # Get download URL
            if attachment_content:
                attachment_text += f"\n\n--- Anexo: {attachment['title']} ---\nURL para Download: {download_url}\n{attachment_content}"  # Added download URL

    full_text = text_with_preserved_structure + attachment_text

    return {"pageContent": full_text, "metadata": metadata}


def load_confluence(base_url, space_key, username, access_token, personal_access_token, cloud=True, process_attachments=True, output_dir="."):  # Added output_dir
    if not valid_base_url(base_url):
        return None, "URL base inválida."

    if not personal_access_token and (not username or not access_token):
        return None, "Você precisa de um token de acesso pessoal (PAT) ou um nome de usuário e token de acesso."

    if not space_key:
        return None, "Você precisa fornecer uma chave de espaço do Confluence."

    st.info(f"-- Trabalhando no Confluence {base_url} --")

    try:
        pages = fetch_all_pages_in_space(base_url, space_key, username, access_token, personal_access_token, cloud)
        documents = [create_document_from_page(page, base_url, space_key, username, access_token, personal_access_token, cloud, process_attachments, output_dir) for page in pages]  # Added output_dir

        documents = [doc for doc in documents if doc is not None]
        return documents, None
    except Exception as e:
        return None, f"Ocorreu um erro ao carregar páginas do Confluence: {e}"


def save_to_file(doc, output_dir, output_format, i):
    """Saves a document to a file in the specified format."""
    title = doc.get("metadata", {}).get("title", f"Page {i}")
    filename = slugify(str(title)) + f"_{i}"

    filepath = os.path.join(output_dir, filename)
    page_content = doc["pageContent"]
    metadata = doc.get('metadata', {})
    file_content = ""

    if output_format == "txt":
        filepath += ".txt"
        file_content = f"Título: {title}\nURL: {metadata.get('url', 'N/A')}\n\n{page_content}"  # Format as plain text
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(file_content)

    elif output_format == "md":
        filepath += ".md"
        file_content = f"# {title}\n[URL]({metadata.get('url', 'N/A')})\n\n{page_content}"  # Format as Markdown
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(file_content)

    elif output_format == "pdf":
        filepath += ".pdf"
        page_content_encoded = page_content.encode('utf-8', 'ignore').decode('latin1')  # Encode to UTF-8, ignore errors, then decode to latin1

        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.cell(200, 10, txt=title, ln=1, align="C")
        pdf.cell(200, 10, txt=f"URL: {metadata.get('url', 'N/A')}", ln=1)
        pdf.multi_cell(0, 5, txt=page_content_encoded)  # Use multi_cell for wrapping text
        pdf.output(filepath)

    elif output_format == "json":
        filepath += ".json"
        file_content = json.dumps({"metadata": metadata, "pageContent": page_content}, indent=4, ensure_ascii=False)  # Ensure_ascii = False
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(file_content)

    return filepath


def save_documents(documents, output_dir, output_format):
    timestamp_str = datetime.datetime.now().strftime("%Y%m%d%H%M")
    output_dir = os.path.join(output_dir, "confluence", timestamp_str)
    os.makedirs(output_dir, exist_ok=True)
    filepaths = []
    for i, doc in enumerate(documents):
        filepath = save_to_file(doc, output_dir, output_format, i)
        filepaths.append(filepath)
    return output_dir, filepaths


def create_zip_archive(directory):
    memory_file = io.BytesIO()

    with zipfile.ZipFile(memory_file, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for root, _, files in os.walk(directory):
            for file in files:
                filepath = os.path.join(root, file)
                zipf.write(filepath, os.path.relpath(filepath, directory))

    memory_file.seek(0)
    return memory_file


# Main execution
if st.button("Carregar Dados do Confluence e Criar ZIP"):
    if not all([base_url, space_key]) or (not access_token and not personal_access_token):
        st.error("Por favor, preencha todos os campos obrigatórios na barra lateral.")
    else:
        timestamp_str = datetime.datetime.now().strftime("%Y%m%d%H%M")
        output_dir = os.path.join("download", timestamp_str)
        os.makedirs(os.path.join(output_dir, "anexos"), exist_ok=True)

        with st.spinner("Carregando dados do Confluence..."):
            documents, error = load_confluence(base_url, space_key, username, access_token, personal_access_token, cloud, process_attachments, output_dir)  # Added output_dir

        if error:
            st.error(error)
        elif not documents:
            st.warning("Nenhuma página encontrada ou ocorreu um erro.")
        else:
            output_dir, filepaths = save_documents(documents, output_dir, output_format)
            zip_buffer = create_zip_archive(output_dir)

            zip_filename = f"confluence_export_{timestamp_str}.zip"

            st.download_button(
                label="Baixar Arquivo ZIP",
                data=zip_buffer,
                file_name=zip_filename,
                mime="application/zip",
            )

            st.success(
                f"Carregadas com sucesso {len(documents)} páginas do Confluence, salvas como arquivos `{output_format}` em: `{output_dir}`. Clique no botão abaixo para baixar o arquivo zip.")
